from __future__ import annotations

import asyncio
import json
import math
import os
import re
from datetime import date, datetime, timedelta
from io import BytesIO
from pathlib import Path
from typing import Any, List

from dotenv import load_dotenv
from fastapi import (
    FastAPI,
    File,
    Form,
    HTTPException,
    Request,
    UploadFile,
)
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from google import genai
from google.genai import types
from pydantic import BaseModel

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


# ---------------------------------------------------------
# Application setup
# ---------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent

load_dotenv(BASE_DIR / ".env")

GEMINI_API_KEY = (
    os.getenv("GEMINI_API_KEY", "")
    or os.getenv("GEMMA_API_KEY", "")
).strip()

GEMINI_MODEL = "gemma-4-26b-a4b-it"

GEMINI_SUPPORT_TIMEOUT_SECONDS = 8
GEMINI_ROADMAP_TIMEOUT_SECONDS = 30
GEMINI_CHAT_TIMEOUT_SECONDS = 20

GEMINI_CLIENT = (
    genai.Client(api_key=GEMINI_API_KEY)
    if GEMINI_API_KEY
    else None
)

app = FastAPI(
    title="ZEN Study Companion",
    version="0.2.0",
)

app.mount(
    "/static",
    StaticFiles(
        directory=BASE_DIR / "static"
    ),
    name="static",
)

templates = Jinja2Templates(
    directory=BASE_DIR / "templates"
)


class LessonRequest(BaseModel):
    day: int
    study_date: str
    title: str
    topics: list[str]
    objectives: list[str]
    estimated_minutes: int


STUDY_LEVEL_LIBRARY: dict[str, dict[str, Any]] = {
    "beginner": {
        "label": "Beginner",
        "summary": "Needs smaller steps, more examples and lower cognitive load.",
        "pace_multiplier": 0.88,
        "readiness_score": 0.35,
    },
    "intermediate": {
        "label": "Intermediate",
        "summary": "Can handle standard pacing with a balanced amount of support.",
        "pace_multiplier": 1.0,
        "readiness_score": 0.6,
    },
    "advanced": {
        "label": "Advanced",
        "summary": "Can move faster and handle a little more stretch in each session.",
        "pace_multiplier": 1.08,
        "readiness_score": 0.8,
    },
}

STUDY_MOOD_LIBRARY: dict[str, dict[str, Any]] = {
    "confident": {
        "label": "Confident",
        "summary": "The learner feels ready to move a bit quicker and handle more challenge.",
        "pace_multiplier": 1.03,
        "readiness_adjustment": 0.08,
        "support_mode": "stretch",
    },
    "steady": {
        "label": "Steady",
        "summary": "The learner feels reasonably calm and can follow a balanced pace.",
        "pace_multiplier": 1.0,
        "readiness_adjustment": 0.0,
        "support_mode": "balanced",
    },
    "unsure": {
        "label": "Unsure",
        "summary": "The learner wants clearer steps and more reassurance.",
        "pace_multiplier": 0.93,
        "readiness_adjustment": -0.07,
        "support_mode": "gentle",
    },
    "anxious": {
        "label": "Anxious",
        "summary": "The learner needs calmer pacing, simpler steps and extra reassurance.",
        "pace_multiplier": 0.86,
        "readiness_adjustment": -0.14,
        "support_mode": "high-support",
    },
}


class GradeRequest(BaseModel):
    lesson: dict[str, Any]
    selected_answers: list[int]


LEARNING_STYLE_LIBRARY: dict[str, dict[str, str]] = {
    "quizzes": {
        "label": "Quizzes",
        "summary": "Finish each study block with short knowledge checks.",
        "day_hint": "End with 3 quick questions to test recall.",
        "objective": "Test yourself on {topic} with a short quiz.",
    },
    "flashcards": {
        "label": "Flashcards",
        "summary": "Convert key ideas into question-and-answer cards.",
        "day_hint": "Make 3 flashcards for the biggest ideas.",
        "objective": "Turn {topic} into flashcard prompts and answers.",
    },
    "notes": {
        "label": "Notes",
        "summary": "Use concise summaries, bullets and clean written recall.",
        "day_hint": "Rewrite the topic as a 3-bullet summary.",
        "objective": "Condense {topic} into a short set of notes.",
    },
}

# ---------------------------------------------------------
# Page routes
# ---------------------------------------------------------


@app.get(
    "/",
    response_class=HTMLResponse,
)
async def index(
    request: Request,
) -> HTMLResponse:
    return templates.TemplateResponse(
        request,
        "index.html",
        {
            "request": request,
        },
    )


@app.get(
    "/health",
)
async def health() -> dict[str, str]:
    return {
        "status": "ok",
        "application": "ZEN",
        "gemma_configured": str(
            GEMINI_CLIENT is not None
        ).lower(),
    }


@app.post(
    "/plan",
    response_class=HTMLResponse,
)
async def create_plan_page(
    request: Request,
    prompt: str = Form(...),
    start_date: str = Form(...),
    deadline: str = Form(...),
    study_minutes_per_day: int = Form(...),
    learning_level: str = Form("intermediate"),
    study_mood: str = Form("steady"),
    learning_styles: List[str] = Form(default_factory=list),
    uploaded_files: List[UploadFile] = File(
        default_factory=list
    ),
) -> HTMLResponse:
    plan_data = await build_plan_response(
        prompt=prompt,
        start_date=start_date,
        deadline=deadline,
        study_minutes_per_day=(
            study_minutes_per_day
        ),
        learning_level=learning_level,
        study_mood=study_mood,
        learning_styles=learning_styles,
        uploaded_files=uploaded_files,
    )

    return templates.TemplateResponse(
        request,
        "plan.html",
        {
            "request": request,
            "plan": plan_data["plan"],
            "file_names": plan_data[
                "file_names"
            ],
            "material_context": plan_data[
                "material_context"
            ],
            "learning_styles": plan_data[
                "learning_styles"
            ],
            "study_minutes_per_day": (
                study_minutes_per_day
            ),
            "start_date": start_date,
            "deadline": deadline,
        },
    )


@app.post(
    "/api/plan",
)
async def generate_plan_api(
    prompt: str = Form(...),
    start_date: str = Form(...),
    deadline: str = Form(...),
    study_minutes_per_day: int = Form(...),
    learning_level: str = Form("intermediate"),
    study_mood: str = Form("steady"),
    learning_styles: List[str] = Form(default_factory=list),
    uploaded_files: List[UploadFile] = File(
        default_factory=list
    ),
) -> JSONResponse:
    plan_data = await build_plan_response(
        prompt=prompt,
        start_date=start_date,
        deadline=deadline,
        study_minutes_per_day=study_minutes_per_day,
        learning_level=learning_level,
        study_mood=study_mood,
        learning_styles=learning_styles,
        uploaded_files=uploaded_files,
    )

    return JSONResponse(plan_data)


# ---------------------------------------------------------
# Chat API
# ---------------------------------------------------------

@app.post(
    "/api/chat",
)
async def chat_api(
    request: Request,
) -> JSONResponse:
    content_type = request.headers.get(
        "content-type",
        "",
    )

    message = ""
    follow_up = ""
    context = ""
    uploaded_files: list[Any] = []

    if "multipart/form-data" in content_type:
        form = await request.form()
        message = str(
            form.get("message", "")
        ).strip()
        follow_up = str(
            form.get("follow_up", "")
        ).strip()
        context = str(
            form.get("context", "")
        ).strip()
        uploaded_files = list(
            form.getlist("uploaded_files")
        )
    else:
        try:
            body = await request.json()
        except Exception as exc:
            raise HTTPException(
                status_code=400,
                detail=(
                    "The request must contain "
                    "valid JSON."
                ),
            ) from exc

        message = str(
            body.get(
                "message",
                "",
            )
        ).strip()

        follow_up = str(
            body.get(
                "follow_up",
                "",
            )
        ).strip()

        context = str(
            body.get(
                "context",
                "",
            )
        ).strip()

        uploaded_files = body.get(
            "uploaded_files",
            [],
        )

    if not message:
        raise HTTPException(
            status_code=400,
            detail="A message is required.",
        )

    chat_data = await build_chat_response(
        message=message,
        follow_up=follow_up,
        context=context,
        uploaded_files=uploaded_files,
    )

    return JSONResponse(
        {
            "reply": chat_data[
                "first_reply"
            ],
            "follow_up_reply": chat_data[
                "follow_up_reply"
            ],
        }
    )


@app.post("/api/lesson")
async def generate_lesson_api(
    lesson_request: LessonRequest,
) -> JSONResponse:
    lesson = await generate_level_with_gemma(
        day=lesson_request.day,
        study_date=lesson_request.study_date,
        title=lesson_request.title,
        topics=lesson_request.topics,
        objectives=lesson_request.objectives,
        estimated_minutes=lesson_request.estimated_minutes,
    )

    return JSONResponse(lesson)


@app.post("/api/grade")
async def grade_lesson_api(
    grade_request: GradeRequest,
) -> JSONResponse:
    quiz = grade_request.lesson.get("quiz", [])
    selected_answers = grade_request.selected_answers

    if not isinstance(quiz, list) or not quiz:
        raise HTTPException(
            status_code=400,
            detail="The lesson does not contain a quiz.",
        )

    results: list[dict[str, Any]] = []
    correct_count = 0

    for index, question in enumerate(quiz):
        selected_answer = (
            selected_answers[index]
            if index < len(selected_answers)
            else -1
        )

        correct_answer = int(
            question.get("correct_answer", 0)
        )

        is_correct = (
            selected_answer == correct_answer
        )

        if is_correct:
            correct_count += 1

        options = question.get("options", [])

        selected_text = (
            options[selected_answer]
            if (
                isinstance(options, list)
                and 0 <= selected_answer < len(options)
            )
            else "No answer selected"
        )

        correct_text = (
            options[correct_answer]
            if (
                isinstance(options, list)
                and 0 <= correct_answer < len(options)
            )
            else ""
        )

        results.append(
            {
                "question": question.get(
                    "question",
                    "",
                ),
                "topic": question.get(
                    "topic",
                    "",
                ),
                "selected_answer": selected_text,
                "correct_answer": correct_text,
                "is_correct": is_correct,
                "explanation": question.get(
                    "explanation",
                    "",
                ),
            }
        )

    total_questions = len(quiz)

    score = round(
        (
            correct_count /
            total_questions
        ) * 100
    )

    weak_topics = list(
        dict.fromkeys(
            result["topic"]
            for result in results
            if (
                not result["is_correct"]
                and result["topic"]
            )
        )
    )

    feedback = await generate_quiz_feedback_with_gemma(
        lesson=grade_request.lesson,
        score=score,
        results=results,
        weak_topics=weak_topics,
    )

    return JSONResponse(
        {
            "score": score,
            "correct_count": correct_count,
            "total_questions": total_questions,
            "results": results,
            "weak_topics": weak_topics,
            "feedback": feedback,
        }
    )

# ---------------------------------------------------------
# Date helpers
# ---------------------------------------------------------


def parse_iso_date(
    value: str,
    field_name: str,
) -> date:
    try:
        return datetime.strptime(
            value,
            "%Y-%m-%d",
        ).date()
    except ValueError as exc:
        raise HTTPException(
            status_code=400,
            detail=(
                f"{field_name} must use "
                "the YYYY-MM-DD format."
            ),
        ) from exc


def validate_date_range(
    start_date: str,
    deadline: str,
) -> tuple[date, date]:
    parsed_start = parse_iso_date(
        start_date,
        "Start date",
    )

    parsed_deadline = parse_iso_date(
        deadline,
        "Deadline",
    )

    if parsed_deadline < parsed_start:
        raise HTTPException(
            status_code=400,
            detail=(
                "The deadline cannot be "
                "before the start date."
            ),
        )

    return (
        parsed_start,
        parsed_deadline,
    )


def get_date_range(
    start: date,
    end: date,
) -> list[date]:
    number_of_days = (
        end - start
    ).days + 1

    return [
        start + timedelta(days=index)
        for index in range(
            number_of_days
        )
    ]


def format_display_date(
    study_date: date,
) -> str:
    return study_date.strftime(
        "%a, %d %B %Y"
    )


def format_short_display_date(
    study_date: date,
) -> str:
    return study_date.strftime(
        "%d %b %Y"
    )


def safe_integer(
    value: Any,
    default: int,
) -> int:
    try:
        return int(value)
    except (
        TypeError,
        ValueError,
    ):
        return default


def safe_float(
    value: Any,
    default: float,
) -> float:
    try:
        return float(value)
    except (
        TypeError,
        ValueError,
    ):
        return default


def normalize_learning_level(
    learning_level: str,
) -> str:
    value = str(learning_level).strip().lower()

    if value in STUDY_LEVEL_LIBRARY:
        return value

    return "intermediate"


def normalize_study_mood(
    study_mood: str,
) -> str:
    value = str(study_mood).strip().lower()

    if value in STUDY_MOOD_LIBRARY:
        return value

    return "steady"


def build_heuristic_support_profile(
    learning_level: str,
    study_mood: str,
    study_minutes_per_day: int,
) -> dict[str, Any]:
    level_key = normalize_learning_level(
        learning_level
    )
    mood_key = normalize_study_mood(study_mood)

    level_profile = STUDY_LEVEL_LIBRARY[level_key]
    mood_profile = STUDY_MOOD_LIBRARY[mood_key]

    pace_multiplier = round(
        max(
            0.75,
            min(
                level_profile["pace_multiplier"]
                * mood_profile["pace_multiplier"],
                1.2,
            ),
        ),
        2,
    )

    weighted_minutes = max(
        10,
        min(
            int(round(study_minutes_per_day * pace_multiplier)),
            180,
        ),
    )

    readiness_score = round(
        max(
            0.05,
            min(
                0.95,
                level_profile["readiness_score"]
                + mood_profile["readiness_adjustment"],
            ),
        ),
        2,
    )

    if readiness_score < 0.4:
        support_mode = "high-support"
    elif readiness_score < 0.7:
        support_mode = "balanced"
    else:
        support_mode = "stretch"

    recommendations = [
        "Use short study blocks with one clear goal.",
        "Start each session with a quick example or recap.",
        "End with a very small check-in so the learner gets an early win.",
    ]

    if support_mode == "high-support":
        recommendations[0] = "Break the work into smaller, reassuring steps."
        recommendations[1] = "Keep explanations simple and repeat the key idea once."

    if support_mode == "stretch":
        recommendations[2] = "Add one stretch question or bonus reflection at the end."

    return {
        "title": "Study readiness profile",
        "level_key": level_key,
        "level_label": level_profile["label"],
        "level_summary": level_profile["summary"],
        "mood_key": mood_key,
        "mood_label": mood_profile["label"],
        "mood_summary": mood_profile["summary"],
        "readiness_score": readiness_score,
        "pace_multiplier": pace_multiplier,
        "weighted_minutes": weighted_minutes,
        "support_mode": support_mode,
        "summary": (
            f"Self-report suggests a {level_profile['label'].lower()} learner who feels {mood_profile['label'].lower()}. "
            f"Use a {support_mode.replace('-', ' ')} pace and aim for about {weighted_minutes} minutes per day."
        ),
        "recommendations": recommendations,
        "main_model_guidance": (
            f"Learner level: {level_profile['label']}. Mood: {mood_profile['label']}. "
            f"Readiness score: {readiness_score}/1. Pace multiplier: {pace_multiplier}. "
            f"Support mode: {support_mode}. Prefer lessons that match a {support_mode.replace('-', ' ')} pace."
        ),
    }


async def build_support_profile(
    prompt: str,
    learning_level: str,
    study_mood: str,
    study_minutes_per_day: int,
    learning_styles: List[str],
    raw_text: str,
    file_names: List[str],
) -> dict[str, Any]:
    heuristic_profile = build_heuristic_support_profile(
        learning_level,
        study_mood,
        study_minutes_per_day,
    )

    if GEMINI_CLIENT is None:
        return heuristic_profile

    source_files = ", ".join(
        file_names) if file_names else "No filenames supplied"

    user_prompt = f"""
LEARNER REQUEST:
{prompt.strip() or 'Help me build a realistic study roadmap.'}

SELF-RATED LEVEL:
{heuristic_profile['level_label']}

SELF-RATED FEELING:
{heuristic_profile['mood_label']}

REQUESTED STUDY TIME PER DAY:
{study_minutes_per_day} minutes

SOURCE FILES:
{source_files}

LEARNING STYLES:
{', '.join(learning_styles) if learning_styles else 'balanced mix'}

COURSE MATERIAL PREVIEW:
{raw_text[:3000]}

Return valid JSON only with these keys:
title, summary, level_label, mood_label, readiness_score, pace_multiplier, weighted_minutes, support_mode, recommendations, main_model_guidance
"""

    try:
        response = await asyncio.wait_for(
            asyncio.to_thread(
                GEMINI_CLIENT.models.generate_content,
                model=GEMINI_MODEL,
                contents=user_prompt,
                config=types.GenerateContentConfig(
                    system_instruction=(
                        "You are a study psychology assistant. "
                        "Use the learner's self-rated level and feeling to estimate the support needed. "
                        "Do not diagnose or mention mental health labels beyond the self-reported mood. "
                        "Return valid JSON only."
                    ),
                    temperature=0.2,
                ),
            ),
            timeout=GEMINI_SUPPORT_TIMEOUT_SECONDS,
        )

        parsed_profile = parse_support_profile_response(
            raw_response=response.text or "",
            fallback=heuristic_profile,
        )

        if parsed_profile is not None:
            return parsed_profile

    except asyncio.TimeoutError:
        print(
            f"Gemma support profile request timed out after {GEMINI_SUPPORT_TIMEOUT_SECONDS} seconds."
        )

    except Exception as exc:
        print(f"Gemma support profile request failed: {exc!r}")

    return heuristic_profile


def parse_support_profile_response(
    raw_response: str,
    fallback: dict[str, Any],
) -> dict[str, Any] | None:
    cleaned = clean_json_response(raw_response)

    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        return None

    if not isinstance(parsed, dict):
        return None

    recommendations = parsed.get("recommendations", [])
    if not isinstance(recommendations, list):
        recommendations = []

    merged = dict(fallback)
    merged.update(
        {
            "title": str(parsed.get("title", fallback["title"])),
            "summary": str(parsed.get("summary", fallback["summary"])),
            "level_label": str(parsed.get("level_label", fallback["level_label"])),
            "mood_label": str(parsed.get("mood_label", fallback["mood_label"])),
            "readiness_score": round(
                max(0.05, min(0.95, safe_float(parsed.get(
                    "readiness_score"), fallback["readiness_score"]))),
                2,
            ),
            "pace_multiplier": round(
                max(0.75, min(1.2, safe_float(parsed.get(
                    "pace_multiplier"), fallback["pace_multiplier"]))),
                2,
            ),
            "weighted_minutes": max(
                10,
                min(
                    safe_integer(parsed.get("weighted_minutes"),
                                 fallback["weighted_minutes"]),
                    180,
                ),
            ),
            "support_mode": str(parsed.get("support_mode", fallback["support_mode"])),
            "recommendations": [str(item) for item in recommendations if str(item).strip()] or fallback["recommendations"],
            "main_model_guidance": str(parsed.get("main_model_guidance", fallback["main_model_guidance"])),
        }
    )

    return merged


# ---------------------------------------------------------
# Main plan-building flow
# ---------------------------------------------------------

async def build_plan_response(
    prompt: str,
    start_date: str,
    deadline: str,
    study_minutes_per_day: int,
    learning_level: str,
    study_mood: str,
    learning_styles: List[str],
    uploaded_files: List[UploadFile],
) -> dict[str, Any]:
    selected_learning_styles = normalize_learning_styles(
        learning_styles
    )

    parsed_start, parsed_deadline = (
        validate_date_range(
            start_date,
            deadline,
        )
    )

    study_minutes_per_day = max(
        10,
        min(
            study_minutes_per_day,
            180,
        ),
    )

    file_names: list[str] = []
    extracted_sections: list[str] = []

    for uploaded_file in uploaded_files:
        if not uploaded_file.filename:
            continue

        extracted_text = (
            await extract_text_from_upload(
                uploaded_file
            )
        )

        if extracted_text.strip():
            file_names.append(
                uploaded_file.filename
            )

            extracted_sections.append(
                f"FILE: "
                f"{uploaded_file.filename}\n"
                f"{extracted_text}"
            )

    raw_text = "\n\n".join(
        extracted_sections
    )

    if not raw_text.strip():
        raw_text = (
            load_demo_text_from_folder()
        )

        material_directory = (
            BASE_DIR / "study_material"
        )

        if material_directory.exists():
            file_names = [
                path.name
                for path in sorted(
                    material_directory.glob(
                        "*"
                    )
                )
                if path.is_file()
            ]

    if not raw_text.strip():
        raw_text = (
            "Introduction\n"
            "Core concepts\n"
            "Worked examples\n"
            "Practice questions\n"
            "Revision"
        )

    support_profile = build_heuristic_support_profile(
        learning_level=learning_level,
        study_mood=study_mood,
        study_minutes_per_day=study_minutes_per_day,
    )

    weighted_study_minutes = support_profile.get(
        "weighted_minutes",
        study_minutes_per_day,
    )

    plan = await build_plan_with_gemma(
        prompt=prompt,
        start_date=parsed_start,
        deadline=parsed_deadline,
        study_minutes=weighted_study_minutes,
        raw_text=raw_text,
        file_names=file_names,
        learning_styles=selected_learning_styles,
        support_profile=support_profile,
    )

    plan = decorate_plan_with_learning_styles(
        plan=plan,
        learning_styles=selected_learning_styles,
    )

    plan["support_profile"] = support_profile

    focus_topics = plan.get(
        "focus_topics",
        ["Core concepts"],
    )

    study_session = (
        build_placeholder_study_session(
            focus_topics,
            selected_learning_styles,
        )
    )

    return {
        "plan": plan,
        "file_names": file_names,
        "material_context": raw_text,
        "learning_styles": selected_learning_styles,
    }


# ---------------------------------------------------------
# Gemma roadmap generation
# ---------------------------------------------------------

async def build_plan_with_gemma(
    prompt: str,
    start_date: date,
    deadline: date,
    study_minutes: int,
    raw_text: str,
    file_names: List[str],
    learning_styles: List[str],
    support_profile: dict[str, Any],
) -> dict[str, Any]:
    placeholder_plan = generate_placeholder_plan(
        prompt=prompt,
        start_date=start_date,
        deadline=deadline,
        study_minutes_per_day=study_minutes,
        raw_text=raw_text,
        file_names=file_names,
        learning_styles=learning_styles,
    )

    if GEMINI_CLIENT is None:
        print(
            "GEMINI_API_KEY is not configured. "
            "Using placeholder roadmap."
        )
        return placeholder_plan

    calendar_dates = get_date_range(
        start_date,
        deadline,
    )

    date_list = "\n".join(
        (
            f"- Day {index}: "
            f"{study_date.isoformat()} "
            f"({study_date.strftime('%A')})"
        )
        for index, study_date in enumerate(
            calendar_dates,
            start=1,
        )
    )

    source_files = (
        ", ".join(file_names)
        if file_names
        else "No filenames supplied"
    )

    support_profile_text = (
        f"Level: {support_profile['level_label']}\n"
        f"Mood: {support_profile['mood_label']}\n"
        f"Readiness score: "
        f"{support_profile['readiness_score']}\n"
        f"Pace multiplier: "
        f"{support_profile['pace_multiplier']}\n"
        f"Weighted minutes: "
        f"{support_profile['weighted_minutes']}\n"
        f"Support mode: "
        f"{support_profile['support_mode']}\n"
        f"Recommendations: "
        f"{', '.join(support_profile['recommendations'])}"
    )

    system_instruction = """
You are ZEN, an expert study-planning assistant.

Create a realistic dated study roadmap using only the supplied
lecture material.

You must identify real concepts, definitions, processes and examples
from the uploaded lecture slides.

Do not use filenames, file extensions, slide numbers or metadata as
course topics.

Do not invent unrelated information.

Each roadmap day must include a concise lesson and exactly three quiz
questions based on what was taught.

Return valid JSON only. Do not include markdown fences or commentary.
"""

    user_prompt = f"""
Create a study roadmap using ONLY the lecture material below.

LEARNER REQUEST:
{prompt.strip() or "Help me understand the uploaded lecture material."}

DATES:
{start_date.isoformat()} to {deadline.isoformat()}

AVAILABLE STUDY TIME:
{study_minutes} minutes per day

LEARNER SUPPORT PROFILE:
{support_profile_text}

LEARNING STYLES:
{", ".join(learning_styles) if learning_styles else "balanced"}

SOURCE FILES:
{source_files}

REQUIRED ROADMAP DAYS:
{date_list}

LECTURE MATERIAL:
--- START OF LECTURE MATERIAL ---
{raw_text[:18000]}
--- END OF LECTURE MATERIAL ---

Return exactly one JSON object using this structure:

{{
  "course_title": "Specific title based on the lecture slides",
  "course_summary": "Brief description of the uploaded material",
  "learner_goal": "The learner's goal",
  "focus_topics": [
    "Specific topic found in the slides"
  ],
  "roadmap": [
    {{
      "day": 1,
      "study_date": "YYYY-MM-DD",
      "title": "Specific lesson title based on the slides",
      "estimated_minutes": {study_minutes},
      "topics": [
        "Specific topic from the lecture material"
      ],
      "objectives": [
        "Specific objective based on the lecture material",
        "Another specific objective"
      ],
      "lesson": {{
        "introduction": "Short introduction",
        "explanation": "Concise explanation based on the slides",
        "example": "Concrete example based on the slides",
        "recap": [
          "Key point one",
          "Key point two",
          "Key point three"
        ]
      }},
      "quiz": [
        {{
          "id": "day-1-q-1",
          "topic": "Topic being tested",
          "question": "Question based on the uploaded material",
          "options": [
            "Option one",
            "Option two",
            "Option three",
            "Option four"
          ],
          "correct_answer": 0,
          "explanation": "Why the correct answer is correct"
        }},
        {{
          "id": "day-1-q-2",
          "topic": "Topic being tested",
          "question": "Question based on the uploaded material",
          "options": [
            "Option one",
            "Option two",
            "Option three",
            "Option four"
          ],
          "correct_answer": 1,
          "explanation": "Why the correct answer is correct"
        }},
        {{
          "id": "day-1-q-3",
          "topic": "Topic being tested",
          "question": "Question based on the uploaded material",
          "options": [
            "Option one",
            "Option two",
            "Option three",
            "Option four"
          ],
          "correct_answer": 2,
          "explanation": "Why the correct answer is correct"
        }}
      ]
    }}
  ]
}}

Rules:

1. Create exactly one roadmap entry for every required date.
2. Keep the roadmap entries in the same order as the required dates.
3. Use only facts and topics found in the lecture material.
4. Never use filenames, extensions or slide numbers as topics.
5. Create exactly three quiz questions per day.
6. Every quiz question must have exactly four options.
7. correct_answer must be an integer from 0 to 3.
8. Keep each lesson short enough to complete within {study_minutes} minutes.
9. The quiz must test ideas taught in that day's lesson.
10. Return JSON only.
"""

    try:
        print(
            "Sending extracted lecture material to Gemma:",
            len(raw_text),
            "characters",
        )

        response = await asyncio.wait_for(
            asyncio.to_thread(
                GEMINI_CLIENT.models.generate_content,
                model=GEMINI_MODEL,
                contents=user_prompt,
                config=types.GenerateContentConfig(
                    system_instruction=system_instruction,
                    temperature=0.1,
                    response_mime_type="application/json",
                ),
            ),
            timeout=GEMINI_ROADMAP_TIMEOUT_SECONDS,
        )

        print(
            "Gemma response received:",
            len(response.text or ""),
            "characters",
        )

        parsed_plan = parse_plan_response(
            raw_response=response.text or "",
            study_minutes=study_minutes,
            prompt=prompt,
            start_date=start_date,
            deadline=deadline,
        )

        if parsed_plan is not None:
            print(
                "Gemma roadmap generated successfully:",
                len(
                    parsed_plan.get(
                        "roadmap",
                        [],
                    )
                ),
                "days",
            )

            return parsed_plan

        print(
            "Gemma returned an invalid roadmap. "
            "Using placeholder roadmap."
        )

    except asyncio.TimeoutError:
        print(
            "Gemma roadmap request timed out after "
            f"{GEMINI_ROADMAP_TIMEOUT_SECONDS} seconds."
        )

    except Exception as exc:
        print(
            f"Gemma roadmap request failed: {exc!r}"
        )

    return placeholder_plan


def clean_json_response(
    raw_response: str,
) -> str:
    cleaned = raw_response.strip()

    cleaned = re.sub(
        r"^```(?:json)?\s*",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )

    cleaned = re.sub(
        r"\s*```$",
        "",
        cleaned,
    )

    return cleaned.strip()


def clean_chat_response_text(
    raw_response: str,
) -> str:
    cleaned_lines: list[str] = []

    for raw_line in str(raw_response or "").splitlines():
        line = raw_line.strip()

        if not line:
            cleaned_lines.append("")
            continue

        line = re.sub(r"^\s*[*•]\s+", "- ", line)
        line = re.sub(r"[*_`]+", "", line)
        line = re.sub(r"\s+", " ", line).strip()

        cleaned_lines.append(line)

    cleaned = "\n".join(cleaned_lines).strip()
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)

    return cleaned


def parse_plan_response(
    raw_response: str,
    study_minutes: int,
    prompt: str,
    start_date: date,
    deadline: date,
) -> dict[str, Any] | None:
    cleaned = clean_json_response(
        raw_response
    )

    try:
        parsed = json.loads(
            cleaned
        )
    except json.JSONDecodeError as exc:
        print(
            "Could not parse Gemma "
            f"JSON: {exc}"
        )

        print(
            "Gemma response preview: "
            f"{raw_response[:500]}"
        )

        return None

    if not isinstance(
        parsed,
        dict,
    ):
        return None

    generated_roadmap = parsed.get(
        "roadmap"
    )

    if not isinstance(
        generated_roadmap,
        list,
    ):
        return None

    calendar_dates = get_date_range(
        start_date,
        deadline,
    )

    valid_days: list[
        dict[str, Any]
    ] = []

    for index, expected_date in enumerate(
        calendar_dates,
        start=1,
    ):
        if (
            index - 1 <
            len(generated_roadmap)
            and isinstance(
                generated_roadmap[
                    index - 1
                ],
                dict,
            )
        ):
            day = generated_roadmap[
                index - 1
            ]
        else:
            day = {}

        topics = day.get(
            "topics",
            [],
        )

        objectives = day.get(
            "objectives",
            [],
        )

        quiz = day.get(
            "quiz",
            [],
        )

        lesson = day.get(
            "lesson",
            {},
        )

        if not isinstance(
            topics,
            list,
        ):
            topics = []

        if not isinstance(
            objectives,
            list,
        ):
            objectives = []

        if not isinstance(
            quiz,
            list,
        ):
            quiz = []

        if not isinstance(
            lesson,
            dict,
        ):
            lesson = {}

        clean_topics = [
            str(topic)
            for topic in topics
            if str(topic).strip()
        ]

        if not clean_topics:
            clean_topics = [
                "Core concepts"
            ]

        clean_objectives = [
            str(objective)
            for objective in objectives
            if str(objective).strip()
        ]

        if not clean_objectives:
            clean_objectives = [
                (
                    "Understand the key "
                    "ideas assigned to "
                    "this date."
                ),
                (
                    "Complete a short "
                    "knowledge check."
                ),
            ]

        valid_questions: list[
            dict[str, Any]
        ] = []

        for question_index, question in enumerate(
            quiz[:3],
            start=1,
        ):
            if not isinstance(
                question,
                dict,
            ):
                continue

            options = question.get(
                "options",
                [],
            )

            if not isinstance(
                options,
                list,
            ):
                options = []

            clean_options = [
                str(option)
                for option in options[:4]
            ]

            while len(
                clean_options
            ) < 4:
                clean_options.append(
                    (
                        "Option "
                        f"{len(clean_options) + 1}"
                    )
                )

            correct_answer = safe_integer(
                question.get(
                    "correct_answer",
                    0,
                ),
                0,
            )

            correct_answer = max(
                0,
                min(
                    correct_answer,
                    3,
                ),
            )

            valid_questions.append(
                {
                    "id": str(
                        question.get(
                            "id",
                            (
                                f"day-{index}-"
                                f"q-{question_index}"
                            ),
                        )
                    ),
                    "topic": str(
                        question.get(
                            "topic",
                            clean_topics[0],
                        )
                    ),
                    "question": str(
                        question.get(
                            "question",
                            (
                                "Which statement "
                                "best matches "
                                "today's topic?"
                            ),
                        )
                    ),
                    "options": clean_options,
                    "correct_answer": (
                        correct_answer
                    ),
                    "explanation": str(
                        question.get(
                            "explanation",
                            "",
                        )
                    ),
                }
            )

        while len(
            valid_questions
        ) < 3:
            question_number = (
                len(valid_questions) + 1
            )

            valid_questions.append(
                {
                    "id": (
                        f"day-{index}-"
                        f"q-{question_number}"
                    ),
                    "topic": (
                        clean_topics[0]
                    ),
                    "question": (
                        "Which statement "
                        "best matches "
                        "today's topic?"
                    ),
                    "options": [
                        "The central idea",
                        "An unrelated detail",
                        "A filename",
                        "A deadline",
                    ],
                    "correct_answer": 0,
                    "explanation": (
                        "The first option "
                        "describes the "
                        "central idea."
                    ),
                }
            )

        recap = lesson.get(
            "recap",
            [],
        )

        if not isinstance(
            recap,
            list,
        ):
            recap = []

        estimated_minutes = safe_integer(
            day.get(
                "estimated_minutes",
                study_minutes,
            ),
            study_minutes,
        )

        estimated_minutes = max(
            10,
            min(
                estimated_minutes,
                180,
            ),
        )

        valid_days.append(
            {
                "day": index,
                "study_date": (
                    expected_date.isoformat()
                ),
                "display_date": (
                    format_display_date(
                        expected_date
                    )
                ),
                "weekday_short": (
                    expected_date.strftime(
                        "%a"
                    )
                ),
                "month_short": (
                    expected_date.strftime(
                        "%b"
                    ).upper()
                ),
                "day_number": (
                    expected_date.strftime(
                        "%d"
                    )
                ),
                "title": str(
                    day.get(
                        "title",
                        (
                            f"Study Day "
                            f"{index}"
                        ),
                    )
                ),
                "status": (
                    "available"
                    if index == 1
                    else "locked"
                ),
                "estimated_minutes": (
                    estimated_minutes
                ),
                "topics": clean_topics,
                "objectives": (
                    clean_objectives
                ),
                "lesson": {
                    "introduction": str(
                        lesson.get(
                            "introduction",
                            (
                                "Today's "
                                "learning session"
                            ),
                        )
                    ),
                    "explanation": str(
                        lesson.get(
                            "explanation",
                            (
                                "Review the "
                                "assigned topics "
                                "and explain them "
                                "in your own words."
                            ),
                        )
                    ),
                    "example": str(
                        lesson.get(
                            "example",
                            (
                                "Find an example "
                                "inside your "
                                "uploaded material."
                            ),
                        )
                    ),
                    "recap": [
                        str(item)
                        for item in recap
                    ],
                },
                "quiz": valid_questions,
            }
        )

    focus_topics = parsed.get(
        "focus_topics",
        [],
    )

    if not isinstance(
        focus_topics,
        list,
    ):
        focus_topics = []

    if not focus_topics:
        for roadmap_day in valid_days:
            focus_topics.extend(
                roadmap_day["topics"]
            )

    focus_topics = list(
        dict.fromkeys(
            str(topic)
            for topic in focus_topics
            if str(topic).strip()
        )
    )[:8]

    if not focus_topics:
        focus_topics = [
            "Core concepts"
        ]

    total_days = len(
        valid_days
    )

    return {
        "course_title": str(
            parsed.get(
                "course_title",
                "ZEN Study Roadmap",
            )
        ),
        "course_summary": str(
            parsed.get(
                "course_summary",
                "Your dated roadmap is ready.",
            )
        ),
        "learner_goal": str(
            parsed.get(
                "learner_goal",
                (
                    prompt.strip()
                    or "Make steady progress."
                ),
            )
        ),
        "start_date": (
            start_date.isoformat()
        ),
        "deadline": (
            deadline.isoformat()
        ),
        "start_display_date": (
            format_short_display_date(
                start_date
            )
        ),
        "deadline_display_date": (
            format_short_display_date(
                deadline
            )
        ),
        "start_month": (
            start_date.strftime(
                "%b"
            ).upper()
        ),
        "start_day_number": (
            start_date.strftime(
                "%d"
            )
        ),
        "minutes_per_day": (
            study_minutes
        ),
        "total_days": total_days,
        "midpoint_day": max(
            1,
            math.ceil(
                total_days / 2
            ),
        ),
        "current_day": 1,
        "streak": 0,
        "completed_days": [],
        "weak_topics": [],
        "focus_topics": focus_topics,
        "roadmap": valid_days,
    }


def normalize_learning_styles(learning_styles: List[str]) -> List[str]:
    selected: list[str] = []
    seen: set[str] = set()

    for raw_style in learning_styles:
        style = str(raw_style).strip().lower()

        if not style:
            continue

        if style not in LEARNING_STYLE_LIBRARY:
            continue

        if style in seen:
            continue

        seen.add(style)
        selected.append(style)

    return selected or ["balanced"]


def build_learning_style_guidance(
    learning_styles: List[str],
) -> dict[str, Any]:
    selected = normalize_learning_styles(learning_styles)

    if selected == ["balanced"]:
        return {
            "title": "Balanced study mode",
            "summary": (
                "The roadmap will mix visuals, recall, summaries and quick checks so the plan stays flexible."
            ),
            "highlights": [
                {
                    "label": "Balanced",
                    "description": (
                        "A mix of short explanations, quick practice and clear recap steps."
                    ),
                }
            ],
            "primary_label": "Balanced",
        }

    highlights: list[dict[str, str]] = []
    for style in selected:
        profile = LEARNING_STYLE_LIBRARY[style]
        highlights.append(
            {
                "label": profile["label"],
                "description": profile["summary"],
            }
        )

    primary_label = ", ".join(
        LEARNING_STYLE_LIBRARY[style]["label"]
        for style in selected
    )

    summary = (
        "This roadmap will emphasise "
        f"{primary_label.lower()} so the next page shows the ideas "
        "in the format you learn best."
    )

    return {
        "title": "Learning style roadmap",
        "summary": summary,
        "highlights": highlights,
        "primary_label": primary_label,
    }


def build_style_hint(
    style_key: str,
    topic: str,
) -> str:
    if style_key == "balanced":
        return (
            f"Balanced cue: review {topic} with a short summary, example and recall check."
        )

    profile = LEARNING_STYLE_LIBRARY[style_key]

    return profile["day_hint"].replace("your topic", topic)


def build_style_objective(
    style_key: str,
    topic: str,
) -> str:
    if style_key == "balanced":
        return (
            f"Understand the main ideas behind {topic}."
        )

    profile = LEARNING_STYLE_LIBRARY[style_key]
    return profile["objective"].format(topic=topic)


def decorate_plan_with_learning_styles(
    plan: dict[str, Any],
    learning_styles: List[str],
) -> dict[str, Any]:
    selected = normalize_learning_styles(
        learning_styles
    )

    roadmap = plan.get(
        "roadmap",
        [],
    )

    if not isinstance(
        roadmap,
        list,
    ):
        roadmap = []

    for day_number, day in enumerate(
        roadmap,
        start=1,
    ):
        if not isinstance(day, dict):
            continue

        style_key = selected[
            (day_number - 1) % len(selected)
        ]

        topics = day.get(
            "topics",
            [],
        )

        if not isinstance(topics, list):
            topics = []

        topic = (
            str(topics[0])
            if topics
            else "the topic"
        )

        day["style_key"] = style_key

        day["style_label"] = (
            "Balanced"
            if style_key == "balanced"
            else LEARNING_STYLE_LIBRARY[
                style_key
            ]["label"]
        )

        day["style_hint"] = build_style_hint(
            style_key,
            topic,
        )

        lesson = day.get(
            "lesson",
            {},
        )

        if isinstance(lesson, dict):
            # Keep Gemma's original explanation.
            # Only add the learning-style hint.
            lesson["style_hint"] = (
                day["style_hint"]
            )

            day["lesson"] = lesson

    plan["roadmap"] = roadmap
    plan["learning_styles"] = selected

    plan["learning_style_guidance"] = (
        build_learning_style_guidance(
            selected
        )
    )

    return plan


# ---------------------------------------------------------
# Gemma chat
# ---------------------------------------------------------

async def build_chat_context(
    context: str = "",
    uploaded_files: list[Any] | None = None,
) -> str:
    material_chunks: list[str] = []

    if context.strip():
        material_chunks.append(
            f"Additional context: {context.strip()}"
        )

    for uploaded_file in uploaded_files or []:
        try:
            if hasattr(uploaded_file, "filename"):
                filename = str(
                    getattr(uploaded_file, "filename", "") or ""
                ).strip()
                content = await extract_text_from_upload(uploaded_file)
            else:
                filename = str(
                    uploaded_file.get("filename", "")
                ).strip()
                content = str(
                    uploaded_file.get("content", "")
                ).strip()
        except Exception:
            continue

        if filename and content:
            material_chunks.append(
                f"File: {filename}\n{content[:12000]}"
            )

    return "\n\n".join(material_chunks)


async def build_chat_response(
    message: str,
    follow_up: str = "",
    context: str = "",
    uploaded_files: list[Any] | None = None,
) -> dict[str, str]:
    if GEMINI_CLIENT is None:
        return {
            "first_reply": (
                "The chat assistant is not available right now because the Gemini API key is missing. "
                "Please ask again once the service is configured."
            ),
            "follow_up_reply": "",
        }

    material_context = await build_chat_context(
        context=context,
        uploaded_files=uploaded_files,
    )

    try:
        system_instruction = (
            "You are ZEN, a supportive and concise study companion. "
            "Answer the learner's question using the supplied course material "
            "when available. If the material is insufficient, say so briefly and "
            "offer a helpful next step."
        )

        chat = await asyncio.to_thread(
            GEMINI_CLIENT.chats.create,
            model=GEMINI_MODEL,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.4,
            ),
        )

        user_prompt = message
        if material_context:
            user_prompt = (
                f"{user_prompt}\n\nMaterial context:\n{material_context}"
            )

        first_response = await asyncio.wait_for(
            asyncio.to_thread(chat.send_message, user_prompt),
            timeout=GEMINI_CHAT_TIMEOUT_SECONDS,
        )
        first_reply = clean_chat_response_text(first_response.text or "")

        follow_up_reply = ""
        if follow_up:
            follow_up_prompt = follow_up
            if material_context:
                follow_up_prompt = (
                    f"{follow_up_prompt}\n\nMaterial context:\n{material_context}"
                )

            second_response = await asyncio.wait_for(
                asyncio.to_thread(chat.send_message, follow_up_prompt),
                timeout=GEMINI_CHAT_TIMEOUT_SECONDS,
            )
            follow_up_reply = clean_chat_response_text(
                second_response.text or "")

        return {
            "first_reply": (
                first_reply
                or "I could not generate a reply."
            ),
            "follow_up_reply": follow_up_reply,
        }

    except asyncio.TimeoutError:
        print(
            f"Gemma chat request timed out after {GEMINI_CHAT_TIMEOUT_SECONDS} seconds."
        )
        return {
            "first_reply": "ZEN is taking longer than expected right now. Please try again in a moment.",
            "follow_up_reply": "",
        }

    except Exception as exc:
        exc_text = str(exc)

        if "RESOURCE_EXHAUSTED" in exc_text or "Quota exceeded" in exc_text:
            return {
                "first_reply": "ZEN is a bit busy right now because the Gemini quota was reached. Please try again in a little while.",
                "follow_up_reply": "",
            }

        print("Gemma chat request failed; using fallback reply.")
        return {
            "first_reply": "I’m having trouble reaching the study assistant right now. Please try again in a moment.",
            "follow_up_reply": "",
        }


async def generate_level_with_gemma(
    day: int,
    study_date: str,
    title: str,
    topics: list[str],
    objectives: list[str],
    estimated_minutes: int,
) -> dict[str, Any]:
    """
    Generate one finishable learning level based strictly on
    the roadmap topics clicked by the learner.
    """

    if GEMINI_CLIENT is None:
        return generate_placeholder_level(
            day=day,
            study_date=study_date,
            title=title,
            topics=topics,
            objectives=objectives,
            estimated_minutes=estimated_minutes,
        )

    safe_minutes = max(
        10,
        min(estimated_minutes, 90),
    )

    system_instruction = """
You are ZEN, an adaptive study teacher.

The learner has clicked one level from a roadmap that was already
generated from their course material.

Create a focused micro-course that can be completed today.

You must teach only the supplied roadmap topics and objectives.
Do not replace them with unrelated subjects.
Do not create an entire semester course.
Do not mention that you are an AI.

Return valid JSON only.
Do not include markdown fences or text outside the JSON object.
"""

    prompt = f"""
ROADMAP LEVEL:
Day {day}

ASSIGNED DATE:
{study_date}

LEVEL TITLE:
{title}

MANDATORY TOPICS:
{json.dumps(topics, ensure_ascii=False)}

MANDATORY OBJECTIVES:
{json.dumps(objectives, ensure_ascii=False)}

AVAILABLE TIME:
{safe_minutes} minutes

Create a complete but concise one-day learning level.

The learner must be able to finish the explanation, activity and quiz
within {safe_minutes} minutes.

Return this exact JSON structure:

{{
  "day": {day},
  "study_date": "{study_date}",
  "title": "string",
  "estimated_minutes": {safe_minutes},
  "topics": ["string"],
  "welcome": "string",
  "sections": [
    {{
      "heading": "string",
      "explanation": "string",
      "key_points": [
        "string",
        "string"
      ],
      "example": "string"
    }}
  ],
  "practice_activity": {{
    "title": "string",
    "instructions": "string",
    "expected_minutes": 5
  }},
  "quick_recap": [
    "string",
    "string",
    "string"
  ],
  "quiz": [
    {{
      "id": "level-{day}-q-1",
      "topic": "string",
      "question": "string",
      "options": [
        "string",
        "string",
        "string",
        "string"
      ],
      "correct_answer": 0,
      "explanation": "string"
    }}
  ]
}}

Rules:

1. Use the supplied topics exactly as the subject of the lesson.
2. Create 2 to 4 short teaching sections.
3. Create exactly 3 multiple-choice quiz questions.
4. Every question must have exactly 4 options.
5. correct_answer must be an integer from 0 to 3.
6. The whole learning level must fit within {safe_minutes} minutes.
7. Use simple language and concrete examples.
8. The quiz must test the explanations actually taught.
"""

    try:
        response = GEMINI_CLIENT.models.generate_content(
            model=GEMINI_MODEL,
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.25,
                response_mime_type="application/json",
            ),
        )

        cleaned = clean_json_response(
            response.text or ""
        )

        parsed = json.loads(cleaned)

        validated = validate_generated_level(
            parsed=parsed,
            day=day,
            study_date=study_date,
            title=title,
            topics=topics,
            estimated_minutes=safe_minutes,
        )

        if validated:
            return validated

    except Exception as exc:
        print(
            f"Gemma level generation failed: {exc}"
        )

    return generate_placeholder_level(
        day=day,
        study_date=study_date,
        title=title,
        topics=topics,
        objectives=objectives,
        estimated_minutes=safe_minutes,
    )


def validate_generated_level(
    parsed: Any,
    day: int,
    study_date: str,
    title: str,
    topics: list[str],
    estimated_minutes: int,
) -> dict[str, Any] | None:
    if not isinstance(parsed, dict):
        return None

    sections = parsed.get("sections", [])
    quiz = parsed.get("quiz", [])
    quick_recap = parsed.get("quick_recap", [])
    activity = parsed.get(
        "practice_activity",
        {},
    )

    if not isinstance(sections, list):
        return None

    if not isinstance(quiz, list):
        return None

    if not isinstance(quick_recap, list):
        quick_recap = []

    if not isinstance(activity, dict):
        activity = {}

    clean_sections = []

    for section in sections[:4]:
        if not isinstance(section, dict):
            continue

        key_points = section.get(
            "key_points",
            [],
        )

        if not isinstance(key_points, list):
            key_points = []

        clean_sections.append(
            {
                "heading": str(
                    section.get(
                        "heading",
                        "Key concept",
                    )
                ),
                "explanation": str(
                    section.get(
                        "explanation",
                        "",
                    )
                ),
                "key_points": [
                    str(point)
                    for point in key_points[:5]
                ],
                "example": str(
                    section.get(
                        "example",
                        "",
                    )
                ),
            }
        )

    clean_quiz = []

    for index, question in enumerate(
        quiz[:3],
        start=1,
    ):
        if not isinstance(question, dict):
            continue

        options = question.get(
            "options",
            [],
        )

        if not isinstance(options, list):
            options = []

        options = [
            str(option)
            for option in options[:4]
        ]

        while len(options) < 4:
            options.append(
                f"Option {len(options) + 1}"
            )

        correct_answer = safe_integer(
            question.get(
                "correct_answer",
                0,
            ),
            0,
        )

        correct_answer = max(
            0,
            min(correct_answer, 3),
        )

        clean_quiz.append(
            {
                "id": str(
                    question.get(
                        "id",
                        f"level-{day}-q-{index}",
                    )
                ),
                "topic": str(
                    question.get(
                        "topic",
                        topics[0]
                        if topics
                        else title,
                    )
                ),
                "question": str(
                    question.get(
                        "question",
                        "Which answer is correct?",
                    )
                ),
                "options": options,
                "correct_answer": correct_answer,
                "explanation": str(
                    question.get(
                        "explanation",
                        "",
                    )
                ),
            }
        )

    if not clean_sections or len(clean_quiz) != 3:
        return None

    return {
        "day": day,
        "study_date": study_date,
        "title": str(
            parsed.get(
                "title",
                title,
            )
        ),
        "estimated_minutes": estimated_minutes,
        "topics": topics,
        "welcome": str(
            parsed.get(
                "welcome",
                (
                    "This level is designed "
                    "to be completed today."
                ),
            )
        ),
        "sections": clean_sections,
        "practice_activity": {
            "title": str(
                activity.get(
                    "title",
                    "Quick practice",
                )
            ),
            "instructions": str(
                activity.get(
                    "instructions",
                    (
                        "Explain the main idea "
                        "in your own words."
                    ),
                )
            ),
            "expected_minutes": max(
                2,
                min(
                    safe_integer(
                        activity.get(
                            "expected_minutes",
                            5,
                        ),
                        5,
                    ),
                    15,
                ),
            ),
        },
        "quick_recap": [
            str(point)
            for point in quick_recap[:5]
        ],
        "quiz": clean_quiz,
    }


async def generate_quiz_feedback_with_gemma(
    lesson: dict[str, Any],
    score: int,
    results: list[dict[str, Any]],
    weak_topics: list[str],
) -> dict[str, Any]:
    if GEMINI_CLIENT is None:
        return {
            "headline": (
                "Level complete"
                if score >= 70
                else "A little review will help"
            ),
            "summary": (
                f"You scored {score}%. "
                "Review the explanations for any "
                "questions you missed."
            ),
            "next_step": (
                "Continue to the next level."
                if score >= 70
                else (
                    "Review your weakest topic "
                    "before moving on."
                )
            ),
        }

    prompt = f"""
You are ZEN, a supportive learning evaluator.

LESSON TITLE:
{lesson.get("title", "")}

LESSON TOPICS:
{json.dumps(lesson.get("topics", []), ensure_ascii=False)}

SCORE:
{score}%

QUESTION RESULTS:
{json.dumps(results, ensure_ascii=False)}

WEAK TOPICS:
{json.dumps(weak_topics, ensure_ascii=False)}

Give concise, encouraging feedback.

Return valid JSON only:

{{
  "headline": "string",
  "summary": "string",
  "strengths": ["string"],
  "weaknesses": ["string"],
  "next_step": "string"
}}
"""

    try:
        response = GEMINI_CLIENT.models.generate_content(
            model=GEMINI_MODEL,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.25,
                response_mime_type="application/json",
            ),
        )

        parsed = json.loads(
            clean_json_response(
                response.text or ""
            )
        )

        if isinstance(parsed, dict):
            return parsed

    except Exception as exc:
        print(
            f"Gemma feedback generation failed: {exc}"
        )

    return {
        "headline": "Level complete",
        "summary": f"You scored {score}%.",
        "strengths": [],
        "weaknesses": weak_topics,
        "next_step": (
            "Review weak topics, then continue."
        ),
    }


def generate_placeholder_level(
    day: int,
    study_date: str,
    title: str,
    topics: list[str],
    objectives: list[str],
    estimated_minutes: int,
) -> dict[str, Any]:
    primary_topic = (
        topics[0]
        if topics
        else title
    )

    return {
        "day": day,
        "study_date": study_date,
        "title": title,
        "estimated_minutes": estimated_minutes,
        "topics": topics,
        "welcome": (
            "This focused level can be "
            "completed in one sitting."
        ),
        "sections": [
            {
                "heading": (
                    f"Understanding {primary_topic}"
                ),
                "explanation": (
                    f"Begin by reviewing the main "
                    f"ideas behind {primary_topic}."
                ),
                "key_points": objectives[:3],
                "example": (
                    "Connect this concept to an "
                    "example from your uploaded notes."
                ),
            }
        ],
        "practice_activity": {
            "title": "Explain it simply",
            "instructions": (
                f"Explain {primary_topic} as though "
                "you were teaching it to a friend."
            ),
            "expected_minutes": 5,
        },
        "quick_recap": objectives[:3],
        "quiz": create_placeholder_quiz(
            day_number=day,
            topic=primary_topic,
            next_topic=(
                topics[1]
                if len(topics) > 1
                else primary_topic
            ),
        ),
    }
# ---------------------------------------------------------
# File extraction
# ---------------------------------------------------------


async def extract_text_from_upload(
    uploaded_file: UploadFile,
) -> str:
    contents = await uploaded_file.read()

    suffix = Path(
        uploaded_file.filename or ""
    ).suffix.lower()

    try:
        if suffix == ".pptx":
            if Presentation is None:
                return ""

            presentation = Presentation(
                BytesIO(contents)
            )

            return extract_text_from_pptx(
                presentation
            )

        if suffix == ".docx":
            if Document is None:
                return ""

            document = Document(
                BytesIO(contents)
            )

            return "\n".join(
                paragraph.text.strip()
                for paragraph
                in document.paragraphs
                if paragraph.text.strip()
            )

        if suffix == ".txt":
            return contents.decode(
                "utf-8",
                errors="ignore",
            )

    except Exception as exc:
        raise HTTPException(
            status_code=400,
            detail=(
                "Could not read "
                f"{uploaded_file.filename}: "
                f"{exc}"
            ),
        ) from exc

    return ""


def extract_text_from_pptx(
    presentation: Any,
) -> str:
    slide_sections: list[str] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            shape_text = getattr(
                shape,
                "text",
                "",
            )

            if (
                shape_text
                and shape_text.strip()
            ):
                slide_parts.append(
                    shape_text.strip()
                )

            if getattr(
                shape,
                "has_table",
                False,
            ):
                for row in shape.table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]

                    if cells:
                        slide_parts.append(
                            " | ".join(cells)
                        )

        if slide_parts:
            slide_sections.append(
                f"SLIDE {slide_number}\n"
                + "\n".join(slide_parts)
            )

    extracted_text = "\n\n".join(
        slide_sections
    )

    print(
        "PowerPoint extraction:",
        len(presentation.slides),
        "slides,",
        len(extracted_text),
        "characters",
    )

    return extracted_text


def load_demo_text_from_folder() -> str:
    material_directory = (
        BASE_DIR / "study_material"
    )

    if not material_directory.exists():
        return ""

    collected: list[str] = []

    for file_path in sorted(
        material_directory.glob("*")
    ):
        suffix = (
            file_path.suffix.lower()
        )

        try:
            if suffix == ".pptx":
                if Presentation is None:
                    continue

                presentation = Presentation(
                    str(file_path)
                )

                text = (
                    extract_text_from_pptx(
                        presentation
                    )
                )

            elif suffix == ".docx":
                if Document is None:
                    continue

                document = Document(
                    str(file_path)
                )

                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph
                    in document.paragraphs
                    if paragraph.text.strip()
                )

            elif suffix == ".txt":
                text = file_path.read_text(
                    encoding="utf-8",
                    errors="ignore",
                )

            else:
                continue

            if text.strip():
                collected.append(
                    f"FILE: "
                    f"{file_path.name}\n"
                    f"{text}"
                )

        except Exception as exc:
            print(
                "Could not read "
                f"{file_path.name}: {exc}"
            )

    return "\n\n".join(
        collected
    )


# ---------------------------------------------------------
# Placeholder dated roadmap
# ---------------------------------------------------------

def generate_placeholder_plan(
    prompt: str,
    start_date: date,
    deadline: date,
    study_minutes_per_day: int,
    raw_text: str,
    file_names: list[str],
    learning_styles: List[str] | None = None,
) -> dict[str, Any]:
    topics = extract_placeholder_topics(
        raw_text
    )
    selected_learning_styles = normalize_learning_styles(
        learning_styles or []
    )

    study_minutes = max(
        10,
        min(
            study_minutes_per_day,
            180,
        ),
    )

    calendar_dates = get_date_range(
        start_date,
        deadline,
    )

    total_days = len(calendar_dates)

    roadmap: list[
        dict[str, Any]
    ] = []

    for index, study_date in enumerate(
        calendar_dates,
        start=1,
    ):
        topic = topics[
            (index - 1) % len(topics)
        ]

        next_topic = topics[
            index % len(topics)
        ]

        is_final_day = (
            index == len(
                calendar_dates
            )
        )

        if is_final_day:
            title = (
                "Final review and "
                "course consolidation"
            )

            day_topics = [
                topic,
                "Final review",
            ]
        else:
            title = (
                f"Understanding {topic}"
                if index > 1
                else (
                    "Getting started "
                    f"with {topic}"
                )
            )

            day_topics = [
                topic,
                next_topic,
            ]

        roadmap.append(
            {
                "day": index,
                "study_date": (
                    study_date.isoformat()
                ),
                "display_date": (
                    format_display_date(
                        study_date
                    )
                ),
                "weekday_short": (
                    study_date.strftime(
                        "%a"
                    )
                ),
                "month_short": (
                    study_date.strftime(
                        "%b"
                    ).upper()
                ),
                "day_number": (
                    study_date.strftime(
                        "%d"
                    )
                ),
                "title": title,
                "status": (
                    "available"
                    if index == 1
                    else "locked"
                ),
                "estimated_minutes": (
                    study_minutes
                ),
                "topics": day_topics,
                "objectives": [
                    build_style_objective(
                        selected_learning_styles[
                            (index - 1)
                            % len(selected_learning_styles)
                        ],
                        topic,
                    ),
                    (
                        f"Connect {topic} to "
                        "an example from the "
                        "material."
                    ),
                    (
                        "Complete a short "
                        "knowledge check."
                    ),
                ],
                "lesson": {
                    "introduction": (
                        f"Introduction to "
                        f"{topic}"
                    ),
                    "explanation": (
                        "Review the key ideas "
                        f"connected to {topic} "
                        "and explain them in "
                        "your own words."
                    ),
                    "example": (
                        "Find one example of "
                        f"{topic} inside your "
                        "uploaded material."
                    ),
                    "recap": [
                        (
                            "State the main "
                            f"meaning of {topic}."
                        ),
                        (
                            f"Explain how {topic} "
                            f"connects to "
                            f"{next_topic}."
                        ),
                        (
                            "Write down one "
                            "question you still "
                            "have."
                        ),
                    ],
                },
                "quiz": [
                    {
                        "id": (
                            f"day-{index}-q-1"
                        ),
                        "topic": topic,
                        "question": (
                            "What is the main idea "
                            f"behind {topic}?"
                        ),
                        "options": [
                            "The central definition",
                            "An unrelated detail",
                            "A file name",
                            "A study deadline",
                        ],
                        "correct_answer": 0,
                        "explanation": (
                            "The first option "
                            "describes the core concept."
                        ),
                    },
                    {
                        "id": (
                            f"day-{index}-q-2"
                        ),
                        "topic": topic,
                        "question": (
                            "What is a useful way "
                            f"to review {topic}?"
                        ),
                        "options": [
                            (
                                "Explain it using "
                                "your own words"
                            ),
                            "Ignore all examples",
                            (
                                "Skip the topic "
                                "completely"
                            ),
                            "Only read the title",
                        ],
                        "correct_answer": 0,
                        "explanation": (
                            "Explaining a topic in "
                            "your own words helps "
                            "reveal gaps."
                        ),
                    },
                    {
                        "id": (
                            f"day-{index}-q-3"
                        ),
                        "topic": next_topic,
                        "question": (
                            f"Why should {next_topic} "
                            "be connected to earlier "
                            "topics?"
                        ),
                        "options": [
                            (
                                "It helps build a "
                                "logical understanding"
                            ),
                            (
                                "It makes the deadline "
                                "longer"
                            ),
                            (
                                "It removes the "
                                "course files"
                            ),
                            (
                                "It changes the "
                                "file format"
                            ),
                        ],
                        "correct_answer": 0,
                        "explanation": (
                            "Connections between "
                            "topics support deeper "
                            "understanding."
                        ),
                    },
                ],
                "style_key": selected_learning_styles[
                    (index - 1)
                    % len(selected_learning_styles)
                ],
                "style_label": (
                    "Balanced"
                    if selected_learning_styles[
                        (index - 1)
                        % len(selected_learning_styles)
                    ] == "balanced"
                    else LEARNING_STYLE_LIBRARY[
                        selected_learning_styles[
                            (index - 1)
                            % len(selected_learning_styles)
                        ]
                    ]["label"]
                ),
                "style_hint": build_style_hint(
                    selected_learning_styles[
                        (index - 1)
                        % len(selected_learning_styles)
                    ],
                    topic,
                ),
            }
        )

    source_label = (
        ", ".join(file_names[:2])
        if file_names
        else "your provided material"
    )

    course_title = (
        topics[0].title()
        if topics
        else "Study Course"
    )

    summary = (
        f"ZEN created a {total_days}-day "
        f"study roadmap using {source_label}."
    )

    study_blocks = [
        {
            "title": (
                f"Day {day['day']}: "
                f"{day['title']}"
            ),
            "focus": ", ".join(
                day["topics"]
            ),
            "time": (
                f"{day['estimated_minutes']} "
                "minutes"
            ),
            "goal": day["objectives"][0],
            "style_label": day["style_label"],
            "style_hint": day["style_hint"],
        }
        for day in roadmap
    ]

    return {
        "course_title": course_title,
        "course_summary": summary,
        "learner_goal": (
            prompt.strip()
            or "Make steady study progress."
        ),
        "deadline": (
            deadline.isoformat()
        ),
        "minutes_per_day": study_minutes,
        "total_days": total_days,
        "midpoint_day": max(
            1,
            math.ceil(total_days / 2),
        ),
        "current_day": 1,
        "streak": 0,
        "completed_days": [],
        "weak_topics": [],
        "focus_topics": topics[:6],
        "roadmap": roadmap,
        "learning_styles": selected_learning_styles,
        "learning_style_guidance": build_learning_style_guidance(
            selected_learning_styles
        ),
        "summary": summary,
        "custom_prompt": (
            prompt.strip()
            or "Keep the plan realistic."
        ),
        "study_blocks": study_blocks,
        "minimum_win": (
            "Spend 10 minutes reviewing "
            f"{topics[0]} and write three "
            "key ideas."
        ),
        "motivation_note": (
            "Start with one manageable session. "
            "Progress matters more than perfection."
        ),
        "daily_time": (
            f"{study_minutes} minutes per day"
        ),
    }


def determine_placeholder_days(
    deadline: str,
) -> int:
    match = re.search(
        r"\d+",
        deadline or "",
    )

    if not match:
        return 4

    requested_days = int(
        match.group()
    )

    return max(
        2,
        min(requested_days, 10),
    )


def create_placeholder_quiz(
    day_number: int,
    topic: str,
    next_topic: str,
) -> list[dict[str, Any]]:
    return [
        {
            "id": (
                f"day-{day_number}-q-1"
            ),
            "topic": topic,
            "question": (
                "What is the main idea "
                f"behind {topic}?"
            ),
            "options": [
                "The central definition",
                "An unrelated detail",
                "A filename",
                "A study deadline",
            ],
            "correct_answer": 0,
            "explanation": (
                "The first option "
                "describes the core "
                "concept."
            ),
        },
        {
            "id": (
                f"day-{day_number}-q-2"
            ),
            "topic": topic,
            "question": (
                "What is a useful way "
                f"to review {topic}?"
            ),
            "options": [
                (
                    "Explain it using "
                    "your own words"
                ),
                "Ignore all examples",
                "Skip the topic",
                "Only read the title",
            ],
            "correct_answer": 0,
            "explanation": (
                "Explaining a topic in "
                "your own words reveals "
                "gaps in understanding."
            ),
        },
        {
            "id": (
                f"day-{day_number}-q-3"
            ),
            "topic": next_topic,
            "question": (
                f"Why should {next_topic} "
                "connect to earlier topics?"
            ),
            "options": [
                (
                    "It builds logical "
                    "understanding"
                ),
                (
                    "It makes the "
                    "deadline longer"
                ),
                (
                    "It removes the "
                    "course files"
                ),
                (
                    "It changes the "
                    "file format"
                ),
            ],
            "correct_answer": 0,
            "explanation": (
                "Connections between "
                "topics support deeper "
                "understanding."
            ),
        },
    ]


def extract_placeholder_topics(
    raw_text: str,
) -> list[str]:
    cleaned_lines: list[str] = []

    for line in raw_text.splitlines():
        cleaned = re.sub(
            r"\s+",
            " ",
            line,
        ).strip()

        cleaned = re.sub(
            r"[^A-Za-z0-9 /&()'-]",
            "",
            cleaned,
        ).strip()

        if len(cleaned) < 4:
            continue

        if len(cleaned) > 80:
            continue

        lowered = cleaned.lower()

        if lowered.startswith(
            (
                "file:",
                "slide ",
                "title:",
            )
        ):
            continue

        cleaned_lines.append(
            cleaned
        )

    topics: list[str] = []
    seen: set[str] = set()

    for line in cleaned_lines:
        lowered = line.lower()

        if lowered in seen:
            continue

        if len(
            line.split()
        ) > 9:
            continue

        seen.add(
            lowered
        )

        topics.append(
            line
        )

        if len(topics) >= 8:
            break

    if not topics:
        return [
            "Core concepts",
            "Key examples",
            "Practice questions",
            "Revision",
        ]

    return topics


def build_day_title(
    day_number: int,
    topic: str,
) -> str:
    if day_number == 1:
        return (
            f"Getting started with {topic}"
        )

    return f"Understanding {topic}"


# ---------------------------------------------------------
# Placeholder study session
# ---------------------------------------------------------

def build_placeholder_study_session(
    topics: list[str],
    learning_styles: list[str] | None = None,
) -> dict[str, Any]:
    first_topic = (
        topics[0]
        if topics
        else "the first topic"
    )

    second_topic = (
        topics[1]
        if len(topics) > 1
        else first_topic
    )
    selected_learning_styles = normalize_learning_styles(
        learning_styles or []
    )
    primary_style = selected_learning_styles[0]
    style_label = (
        "Balanced"
        if primary_style == "balanced"
        else LEARNING_STYLE_LIBRARY[primary_style]["label"]
    )
    style_hint = build_style_hint(
        primary_style,
        first_topic,
    )

    return {
        "title": (
            f"Mini study session: {first_topic}"
        ),
        "explanation": (
            "Begin by identifying the main idea "
            f"behind {first_topic}. Then connect it "
            "to one example from your material. "
            f"Style focus: {style_label}. {style_hint}"
        ),
        "flashcards": [
            {
                "question": (
                    "What is the main idea behind "
                    f"{first_topic}?"
                ),
                "answer": (
                    "Describe the concept briefly "
                    "in your own words."
                ),
            },
            {
                "question": (
                    f"How does {second_topic} "
                    "connect to the course?"
                ),
                "answer": (
                    "Connect it to a definition, "
                    "example or process."
                ),
            },
        ],
        "quiz": {
            "question": (
                "What should you review first "
                f"about {first_topic}?"
            ),
            "answer": (
                "Start with its main definition "
                "and one clear example."
            ),
        },
        "open_question": (
            f"Explain {first_topic} as though "
            "you were teaching it to a friend."
        ),
        "style_label": style_label,
        "style_hint": style_hint,
    }
